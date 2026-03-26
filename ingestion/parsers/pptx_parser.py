"""PPTX parser: extracts text from PowerPoint slides."""
from __future__ import annotations

import io

from pptx import Presentation

from ingestion.parsers.base import BaseParser, ParsedDocument


class PPTXParser(BaseParser):
    supported_extensions = [".pptx"]
    supported_mime_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        prs = Presentation(io.BytesIO(content))
        texts: list[str] = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            slide_texts.append(line)
            if slide_texts:
                texts.append("\n".join(slide_texts))
        full_text = "\n\n".join(texts)
        metadata = {
            "slide_count": len(prs.slides),
            "filename": filename,
        }
        return ParsedDocument(
            text=full_text,
            page_count=len(prs.slides),
            metadata=metadata,
        )
