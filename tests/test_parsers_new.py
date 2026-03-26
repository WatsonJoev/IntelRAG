# tests/test_parsers_new.py
import io
from pptx import Presentation
from pptx.util import Inches


def make_pptx_bytes(slide_text: str = "Hello PPTX World") -> bytes:
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    txBox.text_frame.text = slide_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_parser_extracts_text():
    from ingestion.parsers.pptx_parser import PPTXParser
    parser = PPTXParser()
    result = parser.parse(make_pptx_bytes("Hello PPTX"), "test.pptx")
    assert "Hello PPTX" in result.text
    assert result.page_count >= 1


def test_pptx_parser_supported_extensions():
    from ingestion.parsers.pptx_parser import PPTXParser
    p = PPTXParser()
    assert ".pptx" in p.supported_extensions


def test_html_parser_extracts_text():
    from ingestion.parsers.html_parser import HTMLParser
    html = b"<html><body><h1>Title</h1><p>Body text here.</p></body></html>"
    parser = HTMLParser()
    result = parser.parse(html, "test.html")
    assert "Title" in result.text
    assert "Body text here" in result.text


def test_html_parser_supported_extensions():
    from ingestion.parsers.html_parser import HTMLParser
    p = HTMLParser()
    assert ".html" in p.supported_extensions
    assert ".htm" in p.supported_extensions


def test_json_parser_extracts_text():
    from ingestion.parsers.json_xml_parser import JSONXMLParser
    import json
    data = json.dumps({"title": "Test Doc", "body": "Some content here"}).encode()
    parser = JSONXMLParser()
    result = parser.parse(data, "test.json")
    assert "Test Doc" in result.text
    assert "Some content here" in result.text


def test_xml_parser_extracts_text():
    from ingestion.parsers.json_xml_parser import JSONXMLParser
    xml = b"<root><title>XML Title</title><body>XML Body content</body></root>"
    parser = JSONXMLParser()
    result = parser.parse(xml, "test.xml")
    assert "XML Title" in result.text
    assert "XML Body content" in result.text


def test_json_xml_parser_supported_extensions():
    from ingestion.parsers.json_xml_parser import JSONXMLParser
    p = JSONXMLParser()
    assert ".json" in p.supported_extensions
    assert ".xml" in p.supported_extensions
